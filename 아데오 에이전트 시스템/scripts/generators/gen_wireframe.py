"""
gen_wireframe.py — 화면설계서 PPT 생성기 (DE-03)
Usage: python gen_wireframe.py <data.json> <output.pptx>
"""
import json, sys, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def rgb(hex_str):
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))

# ─── 색상 상수 ────────────────────────────────────────────────────────────────
NAVY    = rgb("1C3557")
BLUE    = rgb("2E74B5")
BLUE_LT = rgb("D6E4F0")
GRAY    = rgb("F5F7FA")
GRAY_D  = rgb("CCCCCC")
WHITE   = rgb("FFFFFF")
BLACK   = rgb("1A1A2E")
RED     = rgb("C00000")
GREEN   = rgb("375623")
ORANGE  = rgb("833C00")

# ─── 슬라이드 크기 (16:9 와이드) ─────────────────────────────────────────────
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ─── 헬퍼: 텍스트박스 추가 ────────────────────────────────────────────────────
def add_textbox(slide, left, top, width, height, text,
                font_size=11, bold=False, color=None,
                bg_color=None, align=PP_ALIGN.LEFT,
                wrap=True, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox


def add_rect(slide, left, top, width, height, bg=None, line=None, line_width=Pt(1)):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_table(slide, left, top, width, rows_data, col_widths_ratio,
              header_bg=NAVY, header_fg=WHITE, row_bg1=GRAY, row_bg2=WHITE,
              font_size=9):
    """간단한 표 추가 (pptx 네이티브 테이블)"""
    num_rows = len(rows_data)
    num_cols = len(rows_data[0])
    row_h = Inches(0.28)
    total_h = row_h * num_rows

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, total_h).table

    # 열 너비 설정
    total_w_emu = width
    for ci, ratio in enumerate(col_widths_ratio):
        table.columns[ci].width = int(total_w_emu * ratio)

    # 행 높이
    for ri in range(num_rows):
        table.rows[ri].height = row_h

    for ri, row in enumerate(rows_data):
        is_header = (ri == 0)
        for ci, val in enumerate(row):
            cell = table.cell(ri, ci)
            cell.text = str(val)
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER if ci == 0 or is_header else PP_ALIGN.LEFT
            run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else tf.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            run.font.name = "맑은 고딕"
            run.font.bold = is_header

            if is_header:
                run.font.color.rgb = header_fg
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                run.font.color.rgb = BLACK
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg1 if ri % 2 == 1 else row_bg2

    return table


# ─── 슬라이드 빌더 ────────────────────────────────────────────────────────────

def build_title_slide(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
    slide.shapes.title  # may not exist on blank layout

    # 배경 — 네이비 상단 절반
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H * 0.55, bg=NAVY)
    add_rect(slide, 0, SLIDE_H * 0.55, SLIDE_W, SLIDE_H * 0.45, bg=WHITE)

    # 로고 텍스트
    add_textbox(slide, Inches(0.5), Inches(0.3), Inches(4), Inches(0.5),
                "ADEO GROUP", font_size=14, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # 메인 제목
    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(11), Inches(1.2),
                d["project_name"], font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.7), Inches(2.8), Inches(8), Inches(0.6),
                "화면설계서  DE-03", font_size=16, bold=False, color=rgb("A0C4E8"), align=PP_ALIGN.LEFT)

    # 구분선
    add_rect(slide, Inches(0.7), SLIDE_H * 0.55 + Inches(0.3),
             Inches(4), Inches(0.04), bg=BLUE)

    # 문서 정보
    info = [
        ("버전", d.get("version", "v1.0")),
        ("작성일", d.get("doc_date", "")),
        ("고객사", d.get("client_name", "")),
        ("작성자", d.get("author", "웹기획팀")),
    ]
    for i, (k, v) in enumerate(info):
        add_textbox(slide, Inches(0.7), SLIDE_H * 0.55 + Inches(0.5 + i * 0.35),
                    Inches(2), Inches(0.3), k, font_size=10, bold=True, color=NAVY)
        add_textbox(slide, Inches(2.5), SLIDE_H * 0.55 + Inches(0.5 + i * 0.35),
                    Inches(5), Inches(0.3), v, font_size=10, bold=False, color=BLACK)


def build_toc_slide(prs, d):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 헤더바
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), bg=NAVY)
    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(10), Inches(0.55),
                "목  차", font_size=18, bold=True, color=WHITE)

    # FO / BO 컬럼 분리
    fo_items = [item for item in d.get("screens", []) if item.get("section") == "FO" and item.get("type") == "Page"]
    bo_items = [item for item in d.get("screens", []) if item.get("section") == "BO" and item.get("type") == "Page"]

    y_start = Inches(1.1)
    col_w   = Inches(5.8)

    for col_idx, (label, items) in enumerate([("FO (Front Office)", fo_items), ("BO (Back Office)", bo_items)]):
        lx = Inches(0.4) + col_idx * Inches(6.7)
        add_rect(slide, lx, y_start, col_w, Inches(0.35), bg=BLUE)
        add_textbox(slide, lx + Inches(0.1), y_start, col_w, Inches(0.35),
                    label, font_size=12, bold=True, color=WHITE)

        for j, item in enumerate(items):
            iy = y_start + Inches(0.45 + j * 0.38)
            bg = rgb("E8F3FF") if j % 2 == 0 else WHITE
            add_rect(slide, lx, iy, col_w, Inches(0.35), bg=bg, line=GRAY_D, line_width=Pt(0.5))
            label_text = f"{item.get('screen_id', '')}  {item.get('screen_name', '')}  {item.get('url', '')}"
            add_textbox(slide, lx + Inches(0.15), iy, col_w - Inches(0.2), Inches(0.35),
                        label_text, font_size=9, color=BLACK)


def build_screen_slide(prs, screen, slide_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ─── 헤더바 ───────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.85), bg=NAVY)
    section_color = BLUE if screen.get("section") == "FO" else rgb("2D6A4F")
    add_rect(slide, 0, 0, Inches(0.18), Inches(0.85), bg=section_color)

    add_textbox(slide, Inches(0.3), Inches(0.08), Inches(9), Inches(0.4),
                f"{screen.get('screen_id', '')}  |  {screen.get('screen_name', '')}",
                font_size=15, bold=True, color=WHITE)
    add_textbox(slide, Inches(0.3), Inches(0.48), Inches(7), Inches(0.32),
                f"URL: {screen.get('url', '')}   |   버전: {screen.get('version', 'v1.0')}   |   관련화면: {screen.get('related', '')}",
                font_size=9, color=rgb("A0C4E8"))
    add_textbox(slide, Inches(11.8), Inches(0.25), Inches(1.3), Inches(0.4),
                f"{slide_num} / {total}", font_size=10, color=rgb("A0C4E8"), align=PP_ALIGN.RIGHT)

    # ─── 와이어프레임 영역 ────────────────────────────────────────────────────
    wf_top  = Inches(0.95)
    wf_left = Inches(0.3)
    wf_w    = Inches(6.8)
    wf_h    = Inches(5.4)

    add_rect(slide, wf_left, wf_top, wf_w, wf_h, bg=GRAY, line=GRAY_D)

    wireframe = screen.get("wireframe", "")
    if wireframe:
        add_textbox(slide, wf_left + Inches(0.1), wf_top + Inches(0.1),
                    wf_w - Inches(0.2), wf_h - Inches(0.2),
                    wireframe, font_size=8, color=BLACK,
                    font_name="Consolas")

    # ─── 기능 설명 테이블 ──────────────────────────────────────────────────────
    tbl_left = Inches(7.3)
    tbl_top  = Inches(0.95)
    tbl_w    = Inches(5.8)

    funcs = screen.get("functions", [])
    if funcs:
        rows = [["No", "컴포넌트", "타입", "기능 설명"]]
        for i, f in enumerate(funcs, 1):
            rows.append([str(i), f.get("component", ""), f.get("type", ""), f.get("function", "")])
        add_table(slide, tbl_left, tbl_top, tbl_w, rows,
                  col_widths_ratio=[0.06, 0.22, 0.14, 0.58])

    # ─── 화면 조건 ────────────────────────────────────────────────────────────
    cond_top = tbl_top + Inches(0.32 * (len(funcs) + 1)) + Inches(0.15)
    conditions = screen.get("conditions", {})
    if conditions:
        add_rect(slide, tbl_left, cond_top, tbl_w, Inches(0.32), bg=BLUE)
        add_textbox(slide, tbl_left + Inches(0.1), cond_top, tbl_w, Inches(0.32),
                    "화면 조건", font_size=10, bold=True, color=WHITE)
        cond_rows = [["조건", "내용"]]
        for k, v in conditions.items():
            cond_rows.append([k, v])
        add_table(slide, tbl_left, cond_top + Inches(0.32), tbl_w, cond_rows,
                  col_widths_ratio=[0.28, 0.72],
                  header_bg=BLUE_LT, header_fg=NAVY,
                  font_size=9)

    # ─── 푸터 ─────────────────────────────────────────────────────────────────
    add_rect(slide, 0, SLIDE_H - Inches(0.28), SLIDE_W, Inches(0.28), bg=GRAY)
    add_textbox(slide, Inches(0.3), SLIDE_H - Inches(0.28), Inches(8), Inches(0.28),
                "ADEO GROUP  |  화면설계서  DE-03  |  Confidential",
                font_size=7, color=rgb("888888"))


def generate(data_path: str, output_path: str):
    with open(data_path, encoding="utf-8") as f:
        d = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs, d)
    build_toc_slide(prs, d)

    page_screens = [s for s in d.get("screens", []) if s.get("type") == "Page"]
    total = len(page_screens)
    for i, screen in enumerate(page_screens, 1):
        build_screen_slide(prs, screen, i, total)

    prs.save(output_path)
    print(f"[gen_wireframe] 저장 완료: {output_path}")


if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python gen_wireframe.py <data.json> <output.pptx>")
        sys.exit(1)
    generate(sys.argv[1], sys.argv[2])
